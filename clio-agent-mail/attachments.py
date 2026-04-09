"""
attachments.py — textextraktion ur e-postbilagor

Stödda format:
  Text:     .txt, .md, .csv
  PDF:      .pdf  (pdfplumber)
  Word:     .docx (python-docx)
  Excel:    .xlsx, .xls (openpyxl)
  PowerPoint: .pptx (python-pptx)
  Bilder:   .png, .gif, .jpg, .jpeg (Claude vision via base64)

Bilder returneras som base64-strängar i ett separat fält —
de skickas till Claude som image-innehållsblock, inte som text.
"""
import base64
import csv
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

TEXT_LIMIT = 8000   # max tecken per bilaga att injicera i prompten


@dataclass
class AttachmentResult:
    filename: str
    filepath: str
    mime_type: str
    text: Optional[str] = None          # extraherad text (None om ej läsbar)
    image_b64: Optional[str] = None     # base64-bild för vision (PNG/JPG/GIF)
    image_media_type: Optional[str] = None
    error: Optional[str] = None


def extract(filepath: str | Path) -> AttachmentResult:
    """
    Extraherar text eller bilddata ur en bilaga.
    Returnerar AttachmentResult med text eller image_b64 ifyllt.
    """
    path = Path(filepath)
    suffix = path.suffix.lower()
    result = AttachmentResult(
        filename=path.name,
        filepath=str(path),
        mime_type=_mime(suffix),
    )

    try:
        if suffix in (".txt", ".md"):
            result.text = _read_text(path)
        elif suffix == ".csv":
            result.text = _read_csv(path)
        elif suffix == ".pdf":
            result.text = _read_pdf(path)
        elif suffix == ".docx":
            result.text = _read_docx(path)
        elif suffix in (".xlsx", ".xls"):
            result.text = _read_excel(path)
        elif suffix == ".pptx":
            result.text = _read_pptx(path)
        elif suffix in (".png", ".jpg", ".jpeg", ".gif"):
            result.image_b64, result.image_media_type = _read_image(path, suffix)
        else:
            result.error = f"Format inte stödt: {suffix}"
    except ImportError as e:
        result.error = f"Saknar paket: {e} — kör: pip install {_missing_pkg(suffix)}"
        logger.warning(f"[bilagor] {result.error}")
    except Exception as e:
        result.error = str(e)
        logger.error(f"[bilagor] Fel vid läsning av {path.name}: {e}")

    return result


# ── Textextraktorer ───────────────────────────────────────────────────────────

def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")[:TEXT_LIMIT]


def _read_csv(path: Path) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(row) for row in reader]
    return "\n".join(rows)[:TEXT_LIMIT]


def _read_pdf(path: Path) -> str:
    import pdfplumber
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
    return "\n\n".join(pages)[:TEXT_LIMIT]


def _read_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)[:TEXT_LIMIT]


def _read_excel(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"[{sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c for c in cells):
                parts.append("\t".join(cells))
    return "\n".join(parts)[:TEXT_LIMIT]


def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            slides.append(f"[Bild {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)[:TEXT_LIMIT]


def _read_image(path: Path, suffix: str):
    data = path.read_bytes()
    b64 = base64.standard_b64encode(data).decode()
    media_map = {
        ".png":  "image/png",
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif":  "image/gif",
    }
    return b64, media_map.get(suffix, "image/png")


# ── Hjälpfunktioner ───────────────────────────────────────────────────────────

def _mime(suffix: str) -> str:
    return {
        ".pdf":  "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".png":  "image/png",
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif":  "image/gif",
        ".csv":  "text/csv",
        ".txt":  "text/plain",
        ".md":   "text/markdown",
    }.get(suffix, "application/octet-stream")


def _missing_pkg(suffix: str) -> str:
    return {
        ".pdf":  "pdfplumber",
        ".docx": "python-docx",
        ".xlsx": "openpyxl",
        ".xls":  "openpyxl",
        ".pptx": "python-pptx",
    }.get(suffix, "okänt paket")
