"""
ingest_minne.py — generisk ingest för *minnet-mappar i Dropbox.

Indexerar PDF/DOCX/PPTX/TXT från en källmapp till valfri Qdrant-collection.
Checksum-baserad: hoppar över oförändrade filer. Körs av cron eller manuellt.

Körning:
  python3 ingest_minne.py --source ~/Dropbox/projekt/Capgemini/Skidförbundet/ssfminnet --collection mem_ssf
  python3 ingest_minne.py --source ~/Dropbox/ftg/AIAB/aiabminnet --collection mem_aiab
  python3 ingest_minne.py --source ... --collection ... --force   # tvinga omindexering
"""
from __future__ import annotations
import argparse, hashlib, uuid
from pathlib import Path
from dotenv import load_dotenv
from openai import OpenAI
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

_here = Path(__file__).parent
load_dotenv(_here / ".env", override=True) or load_dotenv(_here.parent / ".env", override=True)

EMBED_MODEL   = "text-embedding-3-small"
EMBED_DIM     = 1536
CHUNK_WORDS   = 500
OVERLAP_WORDS = 60
SUPPORTED_EXT = {".pdf", ".docx", ".pptx", ".txt"}

# ---------------------------------------------------------------------------
def get_client() -> QdrantClient:
    import os
    host = os.getenv("QDRANT_HOST", "localhost")
    port = int(os.getenv("QDRANT_PORT", "6333"))
    return QdrantClient(host=host, port=port)

def ensure_collection(collection: str) -> None:
    c = get_client()
    existing = {col.name for col in c.get_collections().collections}
    if collection not in existing:
        c.create_collection(
            collection_name=collection,
            vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE),
        )
        print(f"[ingest_minne] Skapade collection: {collection}")
    else:
        print(f"[ingest_minne] Collection finns: {collection}")

# ---------------------------------------------------------------------------
def extract_pdf(path: Path) -> str:
    import pymupdf
    doc = pymupdf.open(str(path))
    return "\n".join(page.get_text() for page in doc)

def extract_docx(path: Path) -> str:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    doc = Document(str(path))
    parts = []
    def iter_blocks(parent):
        for child in parent.element.body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, parent)
            elif child.tag == qn("w:tbl"):
                yield Table(child, parent)
    for block in iter_blocks(doc):
        if isinstance(block, Paragraph):
            if block.text.strip():
                parts.append(block.text.strip())
        elif isinstance(block, Table):
            for row in block.rows:
                seen, cells = set(), []
                for cell in row.cells:
                    cid = id(cell._tc)
                    if cid not in seen:
                        seen.add(cid)
                        t = cell.text.strip()
                        if t:
                            cells.append(t)
                if cells:
                    parts.append(" | ".join(cells))
    return "\n".join(parts)

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text_frame.text.strip() for s in slide.shapes if s.has_text_frame]
        if texts:
            parts.append(f"[Bild {i}] " + " ".join(t for t in texts if t))
    return "\n".join(parts)

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":   return extract_pdf(path)
    if ext == ".docx":  return extract_docx(path)
    if ext == ".pptx":  return extract_pptx(path)
    if ext == ".txt":   return path.read_text(errors="ignore")
    return ""

# ---------------------------------------------------------------------------
def chunk_text(text: str) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i : i + CHUNK_WORDS]))
        i += CHUNK_WORDS - OVERLAP_WORDS
    return [c for c in chunks if len(c.strip()) > 80]

def sha256(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()

def existing_hashes(collection: str, source_file: str) -> set[str]:
    c = get_client()
    results, offset = [], None
    while True:
        page, offset = c.scroll(
            collection_name=collection,
            scroll_filter=None,
            limit=256,
            offset=offset,
            with_payload=True,
        )
        for pt in page:
            if pt.payload.get("source_file") == source_file:
                results.append(pt.payload.get("chunk_hash", ""))
        if offset is None:
            break
    return set(results)

# ---------------------------------------------------------------------------
def ingest_file(path: Path, collection: str, force: bool = False) -> int:
    print(f"\n[ingest_minne] {path.name}")
    text = extract_text(path)
    if not text.strip():
        print("  → ingen text, hoppar över")
        return 0

    chunks    = chunk_text(text)
    known     = set() if force else existing_hashes(collection, path.name)
    oai       = OpenAI()
    qdrant    = get_client()
    new_count = 0

    for idx, chunk in enumerate(chunks):
        h = sha256(chunk)
        if h in known:
            continue
        vec = oai.embeddings.create(input=[chunk], model=EMBED_MODEL).data[0].embedding
        qdrant.upsert(
            collection_name=collection,
            points=[PointStruct(
                id=str(uuid.uuid4()),
                vector=vec,
                payload={
                    "title":       path.stem,
                    "summary":     chunk,
                    "source_file": path.name,
                    "chunk_index": idx,
                    "chunk_hash":  h,
                    "file_type":   path.suffix.lower().lstrip("."),
                    "source":      "minnet",
                },
            )],
        )
        new_count += 1

    print(f"  → {new_count} nya chunks (av {len(chunks)} totalt)")
    return new_count

# ---------------------------------------------------------------------------
def main() -> None:
    parser = argparse.ArgumentParser(description="Indexera *minnet-mapp till Qdrant")
    parser.add_argument("--source",     required=True, type=Path, help="Källmapp (Dropbox/*minnet/)")
    parser.add_argument("--collection", required=True, help="Qdrant collection-namn (t.ex. mem_ssf)")
    parser.add_argument("--force",      action="store_true", help="Tvinga omindexering")
    args = parser.parse_args()

    source = args.source.expanduser().resolve()
    if not source.exists():
        print(f"[ingest_minne] FEL: Källmappen finns inte: {source}")
        raise SystemExit(1)

    ensure_collection(args.collection)

    files = [f for f in sorted(source.rglob("*"))
             if f.suffix.lower() in SUPPORTED_EXT and f.is_file()]
    print(f"\n[ingest_minne] {len(files)} filer i {source.name} → {args.collection}")

    total = sum(ingest_file(f, args.collection, args.force) for f in files)
    print(f"\n[ingest_minne] KLAR — {total} nya chunks i '{args.collection}'")

if __name__ == "__main__":
    main()
