"""
ingest_pmo.py — indexerar SSF PMO-dokument (pdf/docx/pptx) i cap_ssf_pmo.

Körning:
  python3 ingest_pmo.py                        # alla filer i corpus_pmo_ssf/
  python3 ingest_pmo.py --file "foo.pdf"       # enskild fil
  python3 ingest_pmo.py --force                # tvinga omindexering
"""
from __future__ import annotations
import argparse, hashlib, sys, uuid
from pathlib import Path
from dotenv import load_dotenv
from openai import OpenAI
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

_here = Path(__file__).parent
load_dotenv(_here / ".env", override=True) or load_dotenv(_here.parent / ".env", override=True)

COLLECTION   = "cap_ssf_pmo"
EMBED_MODEL  = "text-embedding-3-small"
EMBED_DIM    = 1536
CORPUS_DIR   = _here / "corpus_pmo_ssf"
CHUNK_WORDS  = 500
OVERLAP_WORDS = 60

# ---------------------------------------------------------------------------
def get_client() -> QdrantClient:
    return QdrantClient(host="localhost", port=6333)

def ensure_collection() -> None:
    c = get_client()
    existing = {col.name for col in c.get_collections().collections}
    if COLLECTION not in existing:
        c.create_collection(
            collection_name=COLLECTION,
            vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE),
        )
        print(f"[ingest_pmo] Skapade collection: {COLLECTION}")
    else:
        print(f"[ingest_pmo] Collection finns: {COLLECTION}")

# ---------------------------------------------------------------------------
def extract_pdf(path: Path) -> str:
    import pymupdf
    doc = pymupdf.open(str(path))
    return "\n".join(page.get_text() for page in doc)

def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        if texts:
            parts.append(f"[Slide {i}] " + " ".join(texts))
    return "\n".join(parts)

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":   return extract_pdf(path)
    if ext == ".docx":  return extract_docx(path)
    if ext == ".pptx":  return extract_pptx(path)
    if ext == ".md":    return path.read_text(encoding="utf-8", errors="replace")
    return ""

# ---------------------------------------------------------------------------
def chunk_text(text: str) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunk = words[i : i + CHUNK_WORDS]
        chunks.append(" ".join(chunk))
        i += CHUNK_WORDS - OVERLAP_WORDS
    return [c for c in chunks if len(c.strip()) > 80]

def sha256(text: str) -> str:
    return hashlib.sha256(text.encode()).hexdigest()

# ---------------------------------------------------------------------------
def existing_hashes(source_file: str) -> set[str]:
    c = get_client()
    results, offset = [], None
    while True:
        page, offset = c.scroll(
            collection_name=COLLECTION,
            scroll_filter=None,
            limit=256,
            offset=offset,
            with_payload=True,
        )
        for pt in page:
            if pt.payload.get("source_file") == source_file:
                results.append(pt.payload.get("chunk_hash", ""))
        if offset is None:
            break
    return set(results)

def ingest_file(path: Path, force: bool = False) -> int:
    print(f"\n[ingest_pmo] {path.name}")
    text = extract_text(path)
    if not text.strip():
        print("  → ingen text, hoppar över")
        return 0

    chunks  = chunk_text(text)
    title   = path.stem
    known   = set() if force else existing_hashes(path.name)
    oai     = OpenAI()
    qdrant  = get_client()
    new_count = 0

    for idx, chunk in enumerate(chunks):
        h = sha256(chunk)
        if h in known:
            continue
        vec = oai.embeddings.create(input=[chunk], model=EMBED_MODEL).data[0].embedding
        point = PointStruct(
            id=str(uuid.uuid4()),
            vector=vec,
            payload={
                "title":       title,
                "summary":     chunk,
                "source_file": path.name,
                "chunk_index": idx,
                "chunk_hash":  h,
                "file_type":   path.suffix.lower().lstrip("."),
            },
        )
        qdrant.upsert(collection_name=COLLECTION, points=[point])
        new_count += 1

    print(f"  → {new_count} nya chunks (av {len(chunks)} totalt)")
    return new_count

# ---------------------------------------------------------------------------
def main() -> None:
    parser = argparse.ArgumentParser(description="Indexera SSF PMO-dokument")
    parser.add_argument("--file",  help="Enskild fil att indexera")
    parser.add_argument("--force", action="store_true", help="Tvinga omindexering")
    args = parser.parse_args()

    ensure_collection()

    if args.file:
        p = Path(args.file)
        if not p.is_absolute():
            p = CORPUS_DIR / args.file
        ingest_file(p, force=args.force)
        return

    files = sorted(CORPUS_DIR.rglob("*"))
    files = [f for f in files if f.suffix.lower() in {".pdf", ".docx", ".pptx"} and f.is_file()]
    print(f"\n[ingest_pmo] Hittade {len(files)} filer i {CORPUS_DIR.name}")

    total = 0
    for f in files:
        total += ingest_file(f, force=args.force)
    print(f"\n[ingest_pmo] KLAR — {total} nya chunks totalt i '{COLLECTION}'")

if __name__ == "__main__":
    main()
